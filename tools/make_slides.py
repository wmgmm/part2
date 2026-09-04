#!/usr/bin/env python3
"""Generate Matt Mort's slides for the Part 2 deck, in Matt Hayden's style.

Run from the repo root:  python3 tools/make_slides.py
Output: deliverables/Part2_Mort_slides.pptx

Why a separate file rather than editing the master deck: the master is 9.7 MB of
Matt Hayden's work and editing it in place is not reversible. This produces a
small deck carrying his theme and layout, from which the slides can be copied
into position (they are intended to replace the placeholder slides 26, 27 and 28).

The template is his own file, so fonts, colour scheme and slide layouts are
inherited rather than guessed. Geometry and the panel colour were read off his
slide 26 and are repeated in PANEL / TITLE / BODY / LABEL below.

Slide content is derived from the exercises in src/data/missions.js. If an
exercise changes, change it here too, or the deck and the site drift.
"""

import copy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

TEMPLATE = Path.home() / "Downloads" / "AI in the workplace 2.pptx"
OUT = Path("deliverables/Part2_Mort_slides.pptx")

# Read off his slide 26, in inches.
PANEL = (0.32, 0.41, 3.15, 6.78)
PANEL_FILL = RGBColor(0x77, 0x01, 0x00)
TITLE = (3.66, 0.40, 8.76, 1.45)
LABEL = (0.62, 1.61, 2.35, 0.71)
BODY = (3.94, 2.09, 8.69, 4.20)
TITLE_FONT, TITLE_SIZE = "Aptos Black", Pt(54)   # 72pt is his, but these titles are longer
BODY_FONT, BODY_SIZE = "Aptos", Pt(20)
LABEL_SIZE = Pt(36)

SLIDES = [
    ("Skills", "What a skill actually is", [
        "A named, described set of instructions for one bounded job, saved as a file.",
        "You write the method down once instead of re-explaining it every time you open a chat.",
        "Not a Gem, not an agent. A file you attach, so it works in Copilot and Gemini alike.",
        "Today we hand you three. You are not writing one from scratch in fifteen minutes.",
    ]),
    ("Skills", "Read it before you run it", [
        "Open the file first. Front matter, who it is, when to use it, how to do it, what to hand back.",
        "That legibility is the whole argument: you can audit it, and a colleague can improve it.",
        "It also does not leave the organisation when the person who wrote it does.",
        "Every rule in a good skill says which failure it prevents.",
    ]),
    ("Skills", "Build the training", [
        "Attach two files, the training skill and the sustainability plan, then one short instruction.",
        "Watch what it does first: it names what it found and what it is building from.",
        "A tool that quietly falls back on general knowledge produces training that only looks sourced.",
        "Objectives you can observe. It will refuse \u201cunderstand\u201d and tell you why.",
    ]),
    ("Skills", "The gap note is the payload", [
        "The skill must list everywhere the plan was silent and it left something for a person to fill.",
        "Most commitments have no date, no number and no named owner, so it writes NONE STATED.",
        "That is not the tool failing. A ten-year strategy defers annual targets to the plan underneath it.",
        "It is a list of questions worth putting to whoever owns the policy.",
    ]),
    ("Skills", "Two skills, one output", [
        "Add the house style skill alongside. Voice rules V1 to V6, design rules D1 to D5.",
        "They do not fight: one governs what goes in, the other how it reads and how it looks.",
        "Every edit comes back in a change log naming the rule it came from.",
        "Silent improvement gives you a better draft and no way to review it.",
    ]),
    ("Notebook", "Then make the deck", [
        "Gemini Notebook, two sources: the plan and your finished training module.",
        "Studio, Slide Deck, and tell it to follow the module and add nothing that is not in it.",
        "One generation each. Google publishes no limit for this and it is probably about three a day.",
        "Trace a line on a slide back to the module, to its citation, to the PDF. Three steps.",
    ]),
    ("Local", "Running it on your own machine", [
        "The same kind of model, no cloud, nothing leaving the laptop.",
        "What you gain: confidential data stays put, no monthly quota, no subscription.",
        "What it costs: your own hardware, and a smaller model than the ones you used all morning.",
        "Demo only today. Nothing to install in the room.",
    ]),
    ("Check", "Check it before it goes out", [
        "The third skill verifies and does not rewrite. It hands back a ledger and leaves the fixing to you.",
        "Every claim gets a verdict, and the interesting ones are not the failures.",
        "SUPPORTED WITH CAVEAT and PARTIAL mean true, but you went further than the source did.",
        "This morning that job was a 7,000-character prompt. Now it is a file anyone can run.",
    ]),
    ("Workflow", "Experiments to workflows", [
        "Jisc calls where most teams sit stage two: experimenting and exploring.",
        "Stage three is operational: training people, developing workflows, putting support around them.",
        "A prompt scales to one person. A file scales to a process.",
        "Even with identical files, different people supply different inputs and apply different checks.",
    ]),
    ("Workflow", "Give it a home", [
        "The skills and the workflow record go where the work happens, not in your downloads folder.",
        "One line saying which step of which process. One line naming who owns it.",
        "Log changes the way you would log a change to any other procedure.",
        "The AI does the repeatable part. A person still makes and owns the decision.",
    ]),
]



def blank_out(prs):
    """Remove every slide, keeping layouts, theme and masters."""
    ids = prs.slides._sldIdLst
    for sldId in list(ids):
        prs.part.drop_rel(sldId.rId)
        ids.remove(sldId)


def main():
    if not TEMPLATE.exists():
        raise SystemExit(f"Template not found: {TEMPLATE}")
    prs = Presentation(str(TEMPLATE))

    # Grab the recurring corner logo off his slide 26 before emptying the deck.
    corner = None
    for sh in prs.slides[25].shapes:
        if sh.shape_type == 13 and sh.name == "Picture 6":
            corner = (sh.image.blob, sh.image.ext, sh.left, sh.top, sh.width, sh.height)

    layout = next((l for l in prs.slide_layouts if l.name == "Title and Content"),
                  prs.slide_layouts[1])
    blank_out(prs)

    import io
    for label, title, bullets in SLIDES:
        s = prs.slides.add_slide(layout)
        # Drop the layout's own placeholders; we position everything explicitly.
        for ph in list(s.placeholders):
            ph._element.getparent().remove(ph._element)

        panel = s.shapes.add_shape(5, Inches(PANEL[0]), Inches(PANEL[1]),
                                   Inches(PANEL[2]), Inches(PANEL[3]))  # rounded rect
        panel.fill.solid()
        panel.fill.fore_color.rgb = PANEL_FILL
        panel.line.fill.background()
        panel.shadow.inherit = False
        panel.text_frame.text = ""

        lb = s.shapes.add_textbox(Inches(LABEL[0]), Inches(LABEL[1]),
                                  Inches(LABEL[2]), Inches(LABEL[3]))
        r = lb.text_frame.paragraphs[0].add_run()
        r.text = label
        r.font.size = LABEL_SIZE
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        tb = s.shapes.add_textbox(Inches(TITLE[0]), Inches(TITLE[1]),
                                  Inches(TITLE[2]), Inches(TITLE[3]))
        tb.text_frame.word_wrap = True
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = title
        r.font.name, r.font.size = TITLE_FONT, TITLE_SIZE

        bd = s.shapes.add_textbox(Inches(BODY[0]), Inches(BODY[1]),
                                  Inches(BODY[2]), Inches(BODY[3]))
        bd.text_frame.word_wrap = True
        for i, line in enumerate(bullets):
            p = bd.text_frame.paragraphs[0] if i == 0 else bd.text_frame.add_paragraph()
            r = p.add_run()
            r.text = line
            r.font.name, r.font.size = BODY_FONT, BODY_SIZE
            p.space_after = Pt(14)

        if corner:
            blob, ext, left, top, w, h = corner
            s.shapes.add_picture(io.BytesIO(blob), left, top, w, h)

    OUT.parent.mkdir(exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT} ({len(SLIDES)} slides)")


if __name__ == "__main__":
    main()
